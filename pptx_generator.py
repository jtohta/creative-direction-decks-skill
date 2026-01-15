"""
PowerPoint generation for DJ Brand Guide Generator.
Creates professional 2-slide presentations using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import glob

from color_utils import hex_to_cmyk, hex_to_rgb, is_light_color
from narrative_generator import generate_brand_narrative


def find_uploaded_images() -> list:
    """
    Find images uploaded to the sandbox filesystem via container_upload.

    Files uploaded via container_upload are placed in the sandbox's filesystem.
    This function searches common locations to find them.

    Returns:
        List of paths to found image files, sorted alphabetically for consistent ordering
    """
    # Search patterns for uploaded images in the sandbox
    patterns = [
        "/mnt/user/uploads/*.png",
        "/mnt/user/uploads/*.jpg",
        "/mnt/user/uploads/*.jpeg",
        "/mnt/user/*.png",
        "/mnt/user/*.jpg",
        "/mnt/user/*.jpeg",
        "/uploads/*.png",
        "/uploads/*.jpg",
        "*.png",
        "*.jpg",
        "*.jpeg",
    ]

    images = []
    for pattern in patterns:
        found = glob.glob(pattern)
        images.extend(found)

    # Remove duplicates and sort for consistent ordering
    unique_images = sorted(set(images))

    # Debug: print what we found
    print(f"[DEBUG] Found {len(unique_images)} uploaded images:")
    for img in unique_images:
        print(f"  - {img}")

    return unique_images


def create_brand_guide(dj_input, image_prompts, colors, output_path="brand_guide.pptx", visual_pillars=None):
    """
    Create a complete 3-slide DJ brand guide PowerPoint.

    Args:
        dj_input: Dict containing DJ questionnaire data
        image_prompts: List of dicts with "label", "prompt", and "file_id" keys
        colors: Dict with "primary" and "palette" keys
        output_path: Output file path (default: "brand_guide.pptx")
        visual_pillars: Optional list of pillar dicts with "name" key for Slide 03

    Returns:
        str: Path to the created PowerPoint file
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Get blank slide layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Slide 1: Brand Visual Pillars (if provided)
    if visual_pillars:
        create_visual_pillars_slide(prs, blank_layout, dj_input, visual_pillars)

    # Slide 2: Brand Moodboard
    create_moodboard_slide(prs, blank_layout, dj_input, image_prompts)

    # Slide 3: Color Palette
    create_color_palette_slide(prs, blank_layout, dj_input, colors)

    # Save presentation
    prs.save(output_path)
    return output_path


def create_moodboard_slide(prs, layout, dj_input, image_prompts):
    """
    Create Slide 1: Brand Moodboard with 2x2 image grid and narrative.

    Args:
        prs: Presentation object
        layout: Blank slide layout
        dj_input: DJ questionnaire data
        image_prompts: List of image prompt dicts with label and prompt
    """
    slide = prs.slides.add_slide(layout)

    # Find uploaded images from the sandbox filesystem
    uploaded_images = find_uploaded_images()
    print(f"[DEBUG] Found {len(uploaded_images)} images for {len(image_prompts)} prompts")

    # Title
    dj_name = dj_input['dj_name']
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = f"{dj_name.upper()} - OVERALL BRAND MOODBOARD"
    title_paragraph = title_frame.paragraphs[0]
    title_run = title_paragraph.runs[0]
    title_run.font.name = "Fjalla One"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)

    # Image/prompt boxes - 2-column grid layout
    box_width = Inches(4.25)
    box_height = Inches(2.0)
    start_y = Inches(1.0)
    left_x = Inches(0.5)
    right_x = Inches(5.25)
    gap = Inches(0.3)

    # Loop through all prompts dynamically (supports arbitrary number)
    for i, prompt in enumerate(image_prompts):
        row = i // 2
        col = i % 2
        x = left_x if col == 0 else right_x
        y = start_y + row * (box_height + gap)

        # Add label
        label_box = slide.shapes.add_textbox(x, y, box_width, Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = f"[{prompt['label']}]"
        label_para = label_frame.paragraphs[0]
        label_run = label_para.runs[0]
        label_run.font.name = "Fjalla One"
        label_run.font.size = Pt(10)
        label_run.font.bold = True
        label_run.font.color.rgb = RGBColor(102, 102, 102)

        # Add image or text fallback
        content_y = y + Inches(0.35)
        content_height = box_height - Inches(0.35)

        # Try to use uploaded image from sandbox filesystem
        # Images are added in order via container_upload, so index should match
        if i < len(uploaded_images):
            image_path = uploaded_images[i]
            print(f"[DEBUG] Using uploaded image for prompt {i}: {image_path}")
            try:
                add_image_with_aspect_ratio(
                    slide, image_path, x, content_y, box_width, content_height
                )
            except Exception as e:
                print(f"[ERROR] Failed to add image {image_path}: {e}")
                add_text_fallback(slide, prompt['prompt'], x, content_y, box_width, content_height)
        elif 'path' in prompt and os.path.exists(prompt.get('path', '')):
            # Local file path fallback (for testing)
            add_image_with_aspect_ratio(
                slide, prompt['path'], x, content_y, box_width, content_height
            )
        else:
            # No image available - show text prompt
            print(f"[DEBUG] No image found for prompt {i}, using text fallback")
            add_text_fallback(slide, prompt['prompt'], x, content_y, box_width, content_height)

    # Brand narrative paragraph (positioned below the last row)
    num_rows = (len(image_prompts) + 1) // 2  # Ceiling division
    last_row_y = start_y + (num_rows - 1) * (box_height + gap)
    narrative_y = last_row_y + box_height + gap + Inches(0.3)
    narrative = generate_brand_narrative(dj_input)

    # Narrative title
    narrative_title_box = slide.shapes.add_textbox(
        right_x, narrative_y, box_width, Inches(0.25)
    )
    narrative_title_frame = narrative_title_box.text_frame
    narrative_title_frame.text = "BRAND NARRATIVE"
    narrative_title_para = narrative_title_frame.paragraphs[0]
    narrative_title_run = narrative_title_para.runs[0]
    narrative_title_run.font.name = "Fjalla One"
    narrative_title_run.font.size = Pt(12)
    narrative_title_run.font.bold = True
    narrative_title_run.font.color.rgb = RGBColor(0, 0, 0)

    # Narrative text
    narrative_text_box = slide.shapes.add_textbox(
        right_x, narrative_y + Inches(0.3), box_width, Inches(1.2)
    )
    narrative_text_frame = narrative_text_box.text_frame
    narrative_text_frame.text = narrative
    narrative_text_frame.word_wrap = True
    for paragraph in narrative_text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)


def add_image_with_aspect_ratio(slide, image_path, x, y, box_width, box_height):
    """
    Add image to slide with 16:9 aspect ratio preservation and centering.

    Args:
        slide: Slide object
        image_path: Path to image file
        x, y: Top-left position
        box_width, box_height: Available space dimensions
    """
    image_aspect = 16 / 9
    box_aspect = box_width / box_height

    if box_aspect > image_aspect:
        # Box is wider than 16:9 - fit to height
        img_h = box_height
        img_w = img_h * image_aspect
        img_x = x + (box_width - img_w) / 2  # Center horizontally
        img_y = y
    else:
        # Box is taller than 16:9 - fit to width
        img_w = box_width
        img_h = img_w / image_aspect
        img_x = x
        img_y = y + (box_height - img_h) / 2  # Center vertically

    slide.shapes.add_picture(image_path, img_x, img_y, img_w, img_h)


def add_text_fallback(slide, text, x, y, width, height):
    """
    Add text box as fallback when image is not available.

    Args:
        slide: Slide object
        text: Text content
        x, y: Top-left position
        width, height: Box dimensions
    """
    text_box = slide.shapes.add_textbox(x, y, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    # Set fill and border
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    line = text_box.line
    line.color.rgb = RGBColor(204, 204, 204)
    line.width = Pt(1)

    # Format text
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)


def create_visual_pillars_slide(prs, layout, dj_input, visual_pillars):
    """
    Create Slide 3: Brand Visual Pillars with 4-quadrant layout.

    Displays only pillar names in a 2x2 grid with divider lines.
    Full pillar data (descriptions, characteristics) is used internally
    for image generation but not displayed on this slide.

    Args:
        prs: Presentation object
        layout: Blank slide layout
        dj_input: DJ questionnaire data
        visual_pillars: List of pillar dicts with "name" key
    """
    slide = prs.slides.add_slide(layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "BRAND VISUAL PILLARS"
    title_para = title_frame.paragraphs[0]
    title_run = title_para.runs[0]
    title_run.font.name = "Fjalla One"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)

    # Section label
    label_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), Inches(3), Inches(0.3)
    )
    label_frame = label_box.text_frame
    label_frame.text = "[ VISUAL BRAND PILLARS ]"
    label_para = label_frame.paragraphs[0]
    label_run = label_para.runs[0]
    label_run.font.name = "Fjalla One"
    label_run.font.size = Pt(10)
    label_run.font.bold = True
    label_run.font.color.rgb = RGBColor(102, 102, 102)

    # Description text (right side)
    desc_box = slide.shapes.add_textbox(
        Inches(4.0), Inches(0.9), Inches(5.5), Inches(0.4)
    )
    desc_frame = desc_box.text_frame
    desc_frame.text = "The visual themes and motifs that define the brand's aesthetic direction."
    desc_frame.word_wrap = True
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.name = "Helvetica Neue"
    desc_para.font.size = Pt(9)
    desc_para.font.italic = True
    desc_para.font.color.rgb = RGBColor(102, 102, 102)

    # 4-quadrant grid layout
    grid_top = Inches(1.5)
    grid_left = Inches(0.5)
    grid_width = Inches(9)
    grid_height = Inches(3.5)

    quad_width = grid_width / 2
    quad_height = grid_height / 2
    center_x = grid_left + quad_width
    center_y = grid_top + quad_height

    # Horizontal divider line
    h_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        grid_left, center_y - Pt(0.5), grid_width, Pt(1)
    )
    h_line.fill.solid()
    h_line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    h_line.line.fill.background()

    # Vertical divider line
    v_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        center_x - Pt(0.5), grid_top, Pt(1), grid_height
    )
    v_line.fill.solid()
    v_line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    v_line.line.fill.background()

    # Pillar positions (top-left, top-right, bottom-left, bottom-right)
    positions = [
        (grid_left, grid_top, quad_width, quad_height),                    # Top-left
        (center_x, grid_top, quad_width, quad_height),                     # Top-right
        (grid_left, center_y, quad_width, quad_height),                    # Bottom-left
        (center_x, center_y, quad_width, quad_height),                     # Bottom-right
    ]

    # Add pillar names to each quadrant
    for i, pillar in enumerate(visual_pillars):
        if i >= 4:
            break

        x, y, w, h = positions[i]
        pillar_name = pillar.get('name', pillar) if isinstance(pillar, dict) else str(pillar)

        # Create text box centered in quadrant
        text_box = slide.shapes.add_textbox(x, y, w, h)
        text_frame = text_box.text_frame
        text_frame.text = pillar_name.upper()

        # Center text vertically and horizontally
        text_frame.word_wrap = True
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        para.space_after = Pt(0)

        # Add vertical centering by adjusting margins
        text_frame.margin_top = int(h / 2 - Pt(12))
        text_frame.margin_bottom = Pt(0)
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)

        run = para.runs[0]
        run.font.name = "Fjalla One"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)

    # Footer text
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.2), Inches(7), Inches(0.3)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "These pillars guide all visual decision-making for the brand identity."
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.name = "Helvetica Neue"
    footer_para.font.size = Pt(8)
    footer_para.font.color.rgb = RGBColor(153, 153, 153)

    # Page indicator
    page_box = slide.shapes.add_textbox(
        Inches(9.0), Inches(5.2), Inches(0.5), Inches(0.3)
    )
    page_frame = page_box.text_frame
    page_frame.text = "03"
    page_para = page_frame.paragraphs[0]
    page_para.alignment = PP_ALIGN.RIGHT
    page_run = page_para.runs[0]
    page_run.font.name = "Helvetica Neue"
    page_run.font.size = Pt(10)
    page_run.font.color.rgb = RGBColor(153, 153, 153)


def create_color_palette_slide(prs, layout, dj_input, colors):
    """
    Create Slide 2: Color Palette with primary block and palette bars.

    Args:
        prs: Presentation object
        layout: Blank slide layout
        dj_input: DJ questionnaire data
        colors: Dict with "primary" and "palette" keys
    """
    slide = prs.slides.add_slide(layout)

    # Title (centered)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(8), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "BRAND COLOR PALETTE"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.name = "Fjalla One"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)

    # Primary color - rounded rectangle on left
    primary_hex = colors['primary']['hex']
    primary_rgb = hex_to_rgb(primary_hex)
    primary_cmyk = hex_to_cmyk(primary_hex)

    primary_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(1.15), Inches(5.2), Inches(2.3)
    )
    primary_fill = primary_shape.fill
    primary_fill.solid()
    primary_fill.fore_color.rgb = RGBColor(*primary_rgb)
    primary_shape.line.fill.background()  # No border

    # Primary color text overlay - label
    primary_label_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.45), Inches(4.8), Inches(0.35)
    )
    primary_label_frame = primary_label_box.text_frame
    primary_label_frame.text = "[PRIMARY POP COLOR]"
    primary_label_para = primary_label_frame.paragraphs[0]
    primary_label_run = primary_label_para.runs[0]
    primary_label_run.font.name = "Fjalla One"
    primary_label_run.font.size = Pt(20)
    primary_label_run.font.bold = True
    primary_label_run.font.color.rgb = RGBColor(255, 255, 255)

    # Primary color text overlay - hex + CMYK
    cmyk_text = f"{primary_hex} C: {primary_cmyk['c']}% M: {primary_cmyk['m']}% Y:{primary_cmyk['y']}% K:{primary_cmyk['k']}%"
    primary_info_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.82), Inches(4.8), Inches(0.25)
    )
    primary_info_frame = primary_info_box.text_frame
    primary_info_frame.text = cmyk_text
    primary_info_para = primary_info_frame.paragraphs[0]
    primary_info_run = primary_info_para.runs[0]
    primary_info_run.font.name = "Helvetica Neue"
    primary_info_run.font.size = Pt(12)
    primary_info_run.font.color.rgb = RGBColor(255, 255, 255)

    # Palette colors - stacked rounded rectangles on right
    bar_width = Inches(3.6)
    bar_height = Inches(0.42)
    bar_gap = Inches(0.07)
    start_x = Inches(6.0)
    current_y = Inches(1.15)

    for color_item in colors['palette']:
        color_hex = color_item['hex']
        color_rgb = hex_to_rgb(color_hex)
        cmyk = hex_to_cmyk(color_hex)
        is_light = is_light_color(color_hex)
        text_color = RGBColor(0, 0, 0) if is_light else RGBColor(255, 255, 255)

        # Color bar
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_x, current_y, bar_width, bar_height
        )
        bar_fill = bar_shape.fill
        bar_fill.solid()
        bar_fill.fore_color.rgb = RGBColor(*color_rgb)

        # Border for light colors
        if is_light:
            bar_shape.line.color.rgb = RGBColor(204, 204, 204)
            bar_shape.line.width = Pt(1)
        else:
            bar_shape.line.fill.background()  # No border

        # Text on bar
        bar_text = f"{color_hex} C: {cmyk['c']}% M: {cmyk['m']}% Y:{cmyk['y']}% K:{cmyk['k']}%"
        bar_text_box = slide.shapes.add_textbox(
            start_x + Inches(0.12), current_y + Inches(0.08),
            bar_width - Inches(0.24), bar_height - Inches(0.16)
        )
        bar_text_frame = bar_text_box.text_frame
        bar_text_frame.text = bar_text
        bar_text_frame.vertical_anchor = 1  # Middle vertical alignment
        bar_text_para = bar_text_frame.paragraphs[0]
        bar_text_run = bar_text_para.runs[0]
        bar_text_run.font.name = "Helvetica Neue"
        bar_text_run.font.size = Pt(10)
        bar_text_run.font.color.rgb = text_color

        current_y += bar_height + bar_gap

    # Color description blurb
    description = colors.get('description', '')

    blurb_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(3.6), Inches(5.2), Inches(1.35)
    )
    blurb_frame = blurb_box.text_frame
    blurb_frame.text = description
    blurb_frame.word_wrap = True
    for paragraph in blurb_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(8.5)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)


if __name__ == "__main__":
    # Test with example data
    print("This module is designed to be imported and used within the DJ Brand Guide skill.")
    print("See SKILL.md for usage examples.")
